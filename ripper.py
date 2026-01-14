import os
import io
import tkinter as tk
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import pytesseract
from pytesseract import Output

def rip_slides(path, file_handle, format_type, progress_callback=None, save_images=False, use_ocr=False) -> str:
    if not os.path.exists(path):
        return
    
    print(f"Opening {path}...")
    pres = Presentation(path)

    filename = os.path.basename(path)
    file_name_no_ext = os.path.splitext(filename)[0]
    file_handle.write(get_header(f"Source: {filename}", 1, format_type))

    media_folder_name = f"{file_name_no_ext}_Media"
    media_folder = os.path.join(os.path.dirname(path), media_folder_name)

    if save_images:
        if not os.path.exists(media_folder):
            print("making dir.")
            print(f"{media_folder}")
            os.makedirs(media_folder)

    total_slides = len(pres.slides)
    for i, slide in enumerate(pres.slides):
        slide_number = i + 1
        print(f"Processing Slide {slide_number}...")


        file_handle.write(get_header(f"Slide {slide_number}", 2, format_type))

        img_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                clean_text = shape.text_frame.text
                file_handle.write(f"{clean_text}\n")
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_count += 1

                try:
                    image_blob = shape.image.blob
                    ext = shape.image.ext

                    if save_images:
                        img_filename = f"Slide_{slide_number}_Image_{img_count}.{ext}"
                        img_path = os.path.join(media_folder, img_filename)

                        with open(img_path, "wb") as img_file:
                            img_file.write(image_blob)
                            print(f"saved to {img_path}")

                    if use_ocr:
                        image = Image.open(io.BytesIO(image_blob))

                        extracted_text = extract_text_with_confidence(image)
                        
                        if extracted_text.strip():
                            file_handle.write(f"\n\n[IMAGE READ]:\n{extracted_text}\n")
                
                except Exception as e:
                    print(f"Could not read image on slide {slide_number}: {e}")
        
        if progress_callback:
            progress_callback(slide_number, total_slides)

def get_header(text, level, format_type):
    if format_type == "md":
        return f"\n{'#' * level} {text} \n"
    else:
        return f"\n{'=' * 5} {text} {'=' * 5}\n"

# during Tesseract OCR only return text with high confidence
def extract_text_with_confidence(image, threshold=60):
    data = pytesseract.image_to_data(image, output_type=Output.DICT)

    valid_words = []

    n_boxes = len(data['text'])

    for i in range(n_boxes):
        conf = int(data['conf'][i])

        if conf > threshold:
            word = data['text'][i]

            if word.strip():
                valid_words.append(word)
    
    return " ".join(valid_words)



    # for windows users
    #platform_name = platform.system()

    #if platform_name == "Windows":
    #    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'